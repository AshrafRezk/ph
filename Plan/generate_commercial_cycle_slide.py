#!/usr/bin/env python3
"""Generate Zeta Commercial Cycle slide — clean, spacious, executive layout."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# ── Palette ────────────────────────────────────────────────────────────────
C_PRIMARY = RGBColor(0x1B, 0x6F, 0xD8)
C_PRIMARY_DARK = RGBColor(0x0D, 0x47, 0x8A)
C_PRIMARY_SOFT = RGBColor(0xE8, 0xF2, 0xFC)
C_WAVE = RGBColor(0xC5, 0xDF, 0xF5)
C_BG = RGBColor(0xFF, 0xFF, 0xFF)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_TEXT = RGBColor(0x1A, 0x2B, 0x4A)
C_MUTED = RGBColor(0x6B, 0x7C, 0x93)
C_ACCENT = RGBColor(0x00, 0xA1, 0xE0)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Minimal stage copy — title + role only (details belong in speaker notes)
STAGES = [
    ("1", "See the Business", "C-Level"),
    ("2", "Plan & Configure", "SFE · PM"),
    ("3", "Field Execution", "Rep"),
    ("4", "Report & Roll Up", "DM · SFE"),
    ("5", "Coach & Develop", "DM"),
    ("6", "Decide & Adjust", "C-Level"),
]


def set_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def set_line(shape, color, width_pt=0.75):
    shape.line.color.rgb = color
    shape.line.width = Pt(width_pt)


def textbox(slide, left, top, width, height, text, *, size=12, bold=False, color=C_TEXT, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Segoe UI"
    p.alignment = align
    return box


def add_wave_band(slide):
    """Soft decorative wave — visual only, no labels."""
    y_base = Inches(3.55)
    heights = [Inches(0.55), Inches(0.85), Inches(1.05), Inches(0.85), Inches(0.55)]
    xs = [Inches(0.4), Inches(2.9), Inches(5.4), Inches(7.9), Inches(10.4)]
    ws = [Inches(2.7), Inches(2.7), Inches(2.7), Inches(2.7), Inches(2.7)]

    for x, w, h in zip(xs, ws, heights):
        band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, x, y_base - h / 2, w, h)
        set_fill(band, C_WAVE)
        set_line(band, C_WAVE, 0)


def add_stage_node(slide, index, num, title, role):
    """One clean node: circle + title + role. Plenty of horizontal space."""
    n = len(STAGES)
    margin = Inches(0.55)
    usable = SLIDE_W - margin * 2
    col_w = usable / n
    cx = margin + col_w * index + col_w / 2

    # Alternate vertical position for breathing room (wave feel)
    y_offsets = [Inches(2.05), Inches(2.55), Inches(2.95), Inches(2.55), Inches(2.05), Inches(2.35)]
    cy = y_offsets[index]

    circle_d = Inches(0.72)
    circle = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL,
        cx - circle_d / 2,
        cy,
        circle_d,
        circle_d,
    )
    set_fill(circle, C_PRIMARY)
    set_line(circle, C_PRIMARY, 0)
    tf = circle.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = C_WHITE
    p.font.name = "Segoe UI"
    p.alignment = PP_ALIGN.CENTER

    card_w = Inches(1.85)
    card_h = Inches(0.95)
    card = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        cx - card_w / 2,
        cy + Inches(0.82),
        card_w,
        card_h,
    )
    set_fill(card, C_WHITE)
    set_line(card, C_PRIMARY_SOFT, 1)
    card.adjustments[0] = 0.12

    accent = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        cx - card_w / 2,
        cy + Inches(0.82),
        Inches(0.06),
        card_h,
    )
    set_fill(accent, C_PRIMARY)
    set_line(accent, C_PRIMARY, 0)

    textbox(
        slide,
        cx - card_w / 2 + Inches(0.14),
        cy + Inches(0.86),
        card_w - Inches(0.18),
        Inches(0.42),
        title,
        size=11,
        bold=True,
        color=C_PRIMARY_DARK,
        align=PP_ALIGN.CENTER,
    )
    textbox(
        slide,
        cx - card_w / 2 + Inches(0.14),
        cy + Inches(1.24),
        card_w - Inches(0.18),
        Inches(0.28),
        role,
        size=8.5,
        color=C_MUTED,
        align=PP_ALIGN.CENTER,
    )

    return cx, cy + Inches(0.36)


def add_forward_arrows(slide):
    """Thin arrows between columns."""
    margin = Inches(0.55)
    usable = SLIDE_W - margin * 2
    col_w = usable / len(STAGES)
    y = Inches(2.72)

    for i in range(len(STAGES) - 1):
        x1 = margin + col_w * (i + 1) - Inches(0.08)
        x2 = x1 + Inches(0.16)
        arrow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW, x1, y, Inches(0.2), Inches(0.12))
        set_fill(arrow, C_WAVE)
        set_line(arrow, C_WAVE, 0)


def add_loop_back(slide):
    """Curved return path: stage 6 → stage 1."""
    track = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.75),
        Inches(4.55),
        Inches(11.83),
        Inches(0.52),
    )
    set_fill(track, C_PRIMARY_SOFT)
    set_line(track, C_PRIMARY_SOFT, 0)
    track.adjustments[0] = 0.5

    textbox(
        slide,
        Inches(0.75),
        Inches(4.62),
        Inches(11.83),
        Inches(0.38),
        "↻  Market & field insights feed the next cycle  —  budget · CLM · coaching priorities adjust monthly",
        size=11,
        bold=False,
        color=C_PRIMARY_DARK,
        align=PP_ALIGN.CENTER,
    )


def add_speaker_notes(slide):
    notes = slide.notes_slide.notes_text_frame
    notes.text = """SPEAKER NOTES — Zeta Commercial Cycle

1 · See the Business (C-Level)
   Executive Home · Pharmacy Sales Analytics · Promo Budget · Project Management
   BU coverage, market sell-out trends, spend vs budget, initiative progress

2 · Plan & Configure (SFE · PM)
   Admin Console · territory & product alignment · CLM publish · coaching templates · plan cycles

3 · Field Execution (Rep)
   Field Home · Planner · Visit / Call Report · CLM player · HCP feedback capture

4 · Report & Roll Up (DM · SFE)
   Employee time card · Medical Rep 360 · Team KPI Command Center

5 · Coach & Develop (DM)
   Coaching events · dual rep + manager scoring · competency gaps

6 · Decide & Adjust (C-Level)
   Back to Executive Home · reallocate promo budget · refocus CLM · intensify coaching

Demo order: Executive Home → Pharmacy Sales → Admin Console → Rep flow → Team KPI → Coaching → Executive Home

Audience: CEO/CFO (ROI & market) · PM (CLM & projects) · DM/SFE (KPIs & coaching) · Guru rep (planner & calls)
"""


def build_slide(output_path: Path):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # White canvas
    bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    set_fill(bg, C_BG)
    set_line(bg, C_BG, 0)

    # ── Header (minimal) ──
    textbox(
        slide, Inches(0.7), Inches(0.45), Inches(8.5), Inches(0.55),
        "The Zeta Commercial Cycle",
        size=32, bold=True, color=C_PRIMARY_DARK,
    )
    textbox(
        slide, Inches(0.7), Inches(1.05), Inches(9.0), Inches(0.35),
        "From market insight to field execution — and back",
        size=14, color=C_MUTED,
    )

    textbox(
        slide, Inches(10.6), Inches(0.5), Inches(2.1), Inches(0.3),
        "ZETA PHARMA", size=10, bold=True, color=C_PRIMARY_DARK, align=PP_ALIGN.RIGHT,
    )
    textbox(
        slide, Inches(10.6), Inches(0.78), Inches(2.1), Inches(0.25),
        "salesforce", size=9, color=C_ACCENT, align=PP_ALIGN.RIGHT,
    )

    add_wave_band(slide)

    for i, (num, title, role) in enumerate(STAGES):
        add_stage_node(slide, i, num, title, role)

    add_forward_arrows(slide)
    add_loop_back(slide)

    add_speaker_notes(slide)

    prs.save(str(output_path))
    print(f"Created: {output_path}")


if __name__ == "__main__":
    out = Path(__file__).resolve().parent / "Zeta_Commercial_Cycle_Slide.pptx"
    build_slide(out)
